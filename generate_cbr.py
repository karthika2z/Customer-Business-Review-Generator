#!/usr/bin/env python3
"""
generate_cbr.py – Automated Customer Business Review generator.

Usage:
    python3 generate_cbr.py
    python3 generate_cbr.py --data-dir csv --template "template/CBR Template.pptx"
    python3 generate_cbr.py --output "CBR - Acme - 2026-02-20.pptx"

New template slide layout (0-indexed, 22 slides):
  0  Title
  1  Agenda
  2  Customer Priorities          ← manual (red image)
  3  Product Consumption          ← 1 tile + 2 charts (separate images)
  4  Operational Health           ← 3 tiles + 1 table (separate images)
  5  Business Value               ← text only
  6  Support Requests             ← 4 tiles + 1 requestor table
  7  Software Upgrades            ← 2 tiles + 1 lifecycle table
  8  Aviatrix Use Cases           ← native PPTX tables (editable)
  9  Feature Requests             ← counts text + detail table image
 10  Strategic Engagements        ← text only
 11  ZTMM Alignment               ← manual (banner) + keep Picture 10
 12  Client Success Score         ← text scores; keep OLE Object 1
 13  Current Initiatives          ← manual (red image)
 14  Next Steps                   ← text only
 15  Action Plan                  ← manual (red image)
 16-21  Rubric appendix
"""
import argparse
import logging
import sys
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Optional

log = logging.getLogger(__name__)

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from data_loader import CustomerData
import chart_builder as cb


# ── Defaults ─────────────────────────────────────────────────────────────── #
DEFAULT_DATA_DIR = "csv"
DEFAULT_TEMPLATE = "template/CBR Template.pptx"

# Slides 15–20 (1-indexed) are the scoring rubric appendix → delete
SLIDES_TO_REMOVE_1IDX = set()   # appendix slides are kept in the deck

# Top-of-slide threshold: shapes above this y-position are title/header shapes
# and should not be treated as content/notes placeholders.
TITLE_THRESHOLD = Inches(1.2)   # 1.2"


# ── Aspect-ratio helper ───────────────────────────────────────────────────── #
def _h(embed_w: float, fig_w: float, fig_h: float) -> float:
    """Return the embed height that preserves the figure's aspect ratio."""
    return embed_w * fig_h / fig_w


# ============================================================
# PPTX helpers
# ============================================================

def delete_slide(prs: Presentation, index: int):
    """Delete slide at 0-based index."""
    xml_slides = prs.slides._sldIdLst
    slide_part = prs.slides[index].part
    rId = None
    for rid, rel in prs.part.rels.items():
        if rel.reltype.endswith("/slide") and rel._target is slide_part:
            rId = rid
            break
    if rId:
        prs.part.drop_rel(rId)
    else:
        log.warning("delete_slide: could not find relationship for slide index %d", index)
    del xml_slides[index]


def find_shape(slide, name: str):
    """Return first shape whose name matches exactly."""
    for shape in slide.shapes:
        if shape.name == name:
            return shape
    return None


def replace_in_text_frame(tf, old: str, new: str):
    """Replace all occurrences of old→new across every run in tf."""
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)


def _set_tf_content(tf, header: str, bullets: list):
    """
    Replace text-frame content with bullets only (no separate header line —
    the slide template already shows the section heading).
    The first bullet becomes paragraph 0; remaining bullets are appended.
    """
    tf.word_wrap = True
    if not tf.paragraphs:
        return  # degenerate empty text frame — nothing to populate
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    first = bullets[0] if bullets else ""
    if p0.runs:
        p0.runs[0].text = first
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = first
    for bullet in bullets[1:]:
        np_ = tf.add_paragraph()
        np_.level = 0
        np_.add_run().text = bullet


def _find_notes_shape(slide, keywords: list = None):
    """
    Find the notes/content text box on a slide.

    Priority:
      1. Shape below the title area (top > 1.2") whose stripped text is
         exactly "Notes", "Note", starts with "Notes\\n", or is empty.
      2. Shape below the title area containing any of the keywords.
      3. Shape anywhere (excluding PH idx=0 title) containing any keyword.

    Returns None if nothing found.
    """
    kw = [k.lower() for k in (keywords or [])]

    # Pass 1: dedicated "Notes" placeholder below title area
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.top < TITLE_THRESHOLD:
            continue
        txt = shape.text_frame.text.strip().lower()
        if (txt in ("notes", "note", "notes ", "")
                or txt.startswith("notes\n")
                or txt.startswith("notes ")):
            return shape

    # Pass 2: keyword match below title area
    if kw:
        for term in kw:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if shape.top < TITLE_THRESHOLD:
                    continue
                if term in shape.text_frame.text.strip().lower():
                    return shape

    # Pass 3: keyword match anywhere (skip title placeholder idx=0)
    if kw:
        for term in kw:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                try:
                    if (shape.placeholder_format and
                            shape.placeholder_format.idx == 0):
                        continue
                except Exception:
                    pass
                if term in shape.text_frame.text.strip().lower():
                    return shape

    return None


def _update_notes(slide, header: str, bullets: list, keywords: list = None):
    """Find the notes box and populate it with header + bullets."""
    target = _find_notes_shape(slide, keywords)
    if target:
        _set_tf_content(target.text_frame, header, bullets)


def add_image_inches(slide, img_path: Path,
                     left_in: float, top_in: float,
                     width_in: float, height_in: float):
    """Add a picture to a slide at the given inch coordinates."""
    slide.shapes.add_picture(
        str(img_path),
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )


# ── Tile row helper ───────────────────────────────────────────────────────── #
def _place_tiles(slide, tile_paths: list,
                 left_in: float, top_in: float,
                 total_w: float, gap: float = 0.10):
    """
    Place a list of tile images side-by-side preserving aspect ratio.
    Returns the bottom-most y coordinate used.
    """
    n = len(tile_paths)
    tile_w = (total_w - gap * (n - 1)) / n
    tile_h = _h(tile_w, cb.TILE_FW, cb.TILE_FH)
    x = left_in
    for p in tile_paths:
        add_image_inches(slide, p, x, top_in, tile_w, tile_h)
        x += tile_w + gap
    return top_in + tile_h


# ============================================================
# Per-slide update functions
# ============================================================

def update_slide0_title(slide, data: CustomerData, date_str: str):
    """Slide 1 (index 0) – Title slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, "[Customer Name]", data.customer_name)

    # Shape "Text 5" → "Customer Name | Date"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if "|" in text:
            tf = shape.text_frame
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = f"{data.customer_name} | {date_str}"
                for r in tf.paragraphs[0].runs[1:]:
                    r.text = ""
            break

    # Shape "Text 6" → account owner or "Aviatrix Account Team"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "aviatrix account team" in shape.text_frame.text.lower():
            tf = shape.text_frame
            owner = data.account_owner if data.account_owner else "Aviatrix Account Team"
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = owner
            break


def update_slide1_agenda(slide, data: CustomerData):
    """Slide 2 (index 1) – Agenda: replace [Customer Name]."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, "[Customer Name]", data.customer_name)


def _mark_manual(slide, data: CustomerData, manual_img: Path,
                 short: bool = False):
    """
    Mark a slide as manually completed:
      - Replace [Customer Name] tokens.
      - Add a prominent red 'To Be Completed by Sales Team' image,
        centred on the slide (or bannerised for short=True).
    """
    for shape in slide.shapes:
        if shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, "[Customer Name]", data.customer_name)

    fw = cb.BANNER_FW if short else cb.MANUAL_FW
    fh = cb.BANNER_FH if short else cb.MANUAL_FH

    embed_w = min(fw, 12.5)                         # cap to slide width
    embed_h = _h(embed_w, fw, fh)

    if short:
        # For ZTMM: banner above Picture 10 (which starts at y≈2.963")
        left = (13.333 - embed_w) / 2
        top  = 1.35
    else:
        left = (13.333 - embed_w) / 2
        top  = (7.5   - embed_h) / 2

    add_image_inches(slide, manual_img, left, top, embed_w, embed_h)


def update_slide3_consumption(slide, data: CustomerData,
                               tile_img: Path,
                               stacked_img: Path,
                               spend_img: Path):
    """
    Slide 4 (index 3) – Product Consumption.

    Template: Notes box at (0.544", 1.499") 5.485"×2.214"  [PH idx=24]
    Right content area: left=6.15", avail_width≈6.8"
    Layout (top-to-bottom in the right column):
      • Utilization tile
      • Stacked bar chart
      • Monthly spend chart
    """
    # ── Notes text ──────────────────────────────────────────────────────── #
    insights = []
    if not data.df_consumption.empty:
        df = data.df_consumption.dropna(subset=["Account Name"])
        df = df[df["Account Name"].str.strip() != ""]
        fn = df[df["Parameters"] == "billing_firenet_firewall_inst"].copy()
        if len(fn) >= 2:
            fn["Date Month"] = pd.to_datetime(fn["Date Month"])
            fn = fn.sort_values("Date Month")
            if fn.tail(6)["Consumption MRR"].mean() < fn.head(6)["Consumption MRR"].mean() * 0.8:
                insights.append("Firenet usage declined in transit layer")
    if data.copilot_deployed == 0:
        insights.append("Copilot not deployed — enable for visibility & control")
    pct_raw = data.consumption_pct
    pct = round(pct_raw if pct_raw > 1 else pct_raw * 100, 1)
    if pct > 0:
        insights.append(f"Utilization: {pct}% of contracted capacity")
    if not insights:
        insights.append("Review consumption trends in chart")
    _update_notes(slide, "Product Consumption", insights,
                  keywords=["consumption", "utilization"])

    # ── Images (right column) ──────────────────────────────────────────── #
    LEFT, AVAIL_W, TOP_START = 6.15, 6.8, 1.50

    # Utilization tile – centre it horizontally in the column
    tile_w = 3.20
    tile_h = _h(tile_w, cb.TILE_FW, cb.TILE_FH)
    tile_left = LEFT + (AVAIL_W - tile_w) / 2
    add_image_inches(slide, tile_img, tile_left, TOP_START, tile_w, tile_h)

    y = TOP_START + tile_h + 0.12

    # Stacked bar
    sb_w = AVAIL_W
    sb_h = _h(sb_w, cb.STACKED_FW, cb.STACKED_FH)
    add_image_inches(slide, stacked_img, LEFT, y, sb_w, sb_h)
    y += sb_h + 0.12

    # Monthly spend
    sp_w = AVAIL_W
    sp_h = _h(sp_w, cb.SPEND_FW, cb.SPEND_FH)
    add_image_inches(slide, spend_img, LEFT, y, sp_w, sp_h)


def update_slide4_operational(slide, data: CustomerData,
                               tile1: Path, tile2: Path, tile3: Path,
                               table_img: Path):
    """
    Slide 5 (index 4) – Operational Health.

    Template: Notes box at (0.529", 1.499") 4.091"×4.591"
    Right content area: left=4.72", avail_width≈8.2"
    Layout: row of 3 tiles, then controller table below
    """
    insights = []
    if data.unsupported_controllers > 0:
        insights.append(f"{data.unsupported_controllers} controller(s) on unsupported release")
    if data.copilot_deployed == 0:
        insights.append("CoPilot not deployed for network visibility")
    if not insights:
        insights.append("All controllers on supported releases")
    _update_notes(slide, "Operational Readiness", insights,
                  keywords=["operational", "readiness", "health"])

    LEFT, AVAIL_W, TOP_START = 4.72, 8.20, 1.50

    bottom_of_tiles = _place_tiles(slide, [tile1, tile2, tile3],
                                   LEFT, TOP_START, AVAIL_W)
    y = bottom_of_tiles + 0.15

    # Controller table
    ct_w = AVAIL_W
    ct_h = _h(ct_w, cb.CTRL_FW, cb.CTRL_FH)
    add_image_inches(slide, table_img, LEFT, y, ct_w, ct_h)


def update_slide5_business_value(slide, data: CustomerData):
    """Slide 6 (index 5) – Business Value Realized."""
    if data.bv_status in ("Not Initiated", "", None):
        bullets = ["Business Value assessment not yet initiated."]
    else:
        bullets = [f"Status: {data.bv_status}"]
    if data.cbr_status:
        bullets.append(f"CBR Status: {data.cbr_status}")
    _update_notes(slide, "Business Value", bullets,
                  keywords=["business value"])


def update_slide6_reliability(slide, data: CustomerData,
                               tile1: Path, tile2: Path,
                               tile3: Path, tile4: Path,
                               requestor_img: Path):
    """
    Slide 7 (index 6) – Support Requests.

    Template: Notes box at (0.544", 1.536") 3.349"×4.821"
    Right content area: left=3.97", avail_width≈8.9"
    Layout: row of 4 tiles, then requestor table below
    """
    if data.p1_tickets > 0:
        bullets = [f"{data.p1_tickets} P1 ticket(s) — requires immediate attention"]
    elif data.total_tickets > 0:
        bullets = [f"{data.total_tickets} ticket(s) in past 12 months"]
    else:
        bullets = ["No tickets opened in the past 12 months"]
    _update_notes(slide, "Support Request", bullets,
                  keywords=["support", "reliability", "ticket"])

    LEFT, AVAIL_W, TOP_START = 3.97, 8.90, 1.55

    bottom_of_tiles = _place_tiles(slide, [tile1, tile2, tile3, tile4],
                                   LEFT, TOP_START, AVAIL_W)
    y = bottom_of_tiles + 0.15

    # Requestor table
    rq_w = AVAIL_W
    rq_h = _h(rq_w, cb.REQ_FW, cb.REQ_FH)
    add_image_inches(slide, requestor_img, LEFT, y, rq_w, rq_h)


def update_slide7_software_upgrades(slide, data: CustomerData,
                                     tile1: Path, tile2: Path,
                                     lifecycle_img: Path):
    """
    Slide 8 (index 7) – Software Upgrades.

    Template: Notes box at (0.529", 1.499") 3.780"×4.591"
    Right content area: left=4.40", avail_width≈8.6"
    Layout: row of 2 tiles, then lifecycle table below
    """
    bullets = []
    if data.upgrade_tickets_count > 0:
        bullets.append(f"{data.upgrade_tickets_count} upgrade ticket(s) currently open")
    else:
        bullets.append("No open upgrade tickets")
    bullets.append("Recommended: stay on latest supported release (N or N-1)")
    if data.gateway_sum > 0:
        bullets.append(f"{data.gateway_sum} total gateways in environment")
    _update_notes(slide, "Software Upgrades", bullets,
                  keywords=["software upgrade", "upgrades"])

    LEFT, AVAIL_W, TOP_START = 4.40, 8.60, 1.50

    bottom_of_tiles = _place_tiles(slide, [tile1, tile2],
                                   LEFT, TOP_START, AVAIL_W)
    y = bottom_of_tiles + 0.15

    # Release lifecycle table
    rl_w = AVAIL_W
    rl_h = _h(rl_w, cb.REL_FW, cb.REL_FH)
    add_image_inches(slide, lifecycle_img, LEFT, y, rl_w, rl_h)


def update_slide8_usecases(slide, data: CustomerData):
    """
    Slide 9 (index 8) – Aviatrix Use Cases.
    Builds two native PPTX tables (Features Enabled + Use Cases) so the
    content remains editable after generation.
    """
    _NAVY   = RGBColor(0x1E, 0x2D, 0x3D)
    _WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
    _LBLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # light blue
    _GREEN  = RGBColor(0x27, 0xAE, 0x60)
    _RED    = RGBColor(0xE7, 0x4C, 0x3C)
    _ROW_BG = [_WHITE, _LBLUE]

    features = [
        ("Copilot Enabled",            data.copilot_enabled_label == "Y"),
        ("DCF Enabled",                data.dcf_enabled == "Y"),
        ("Firenet Enabled",            data.firenet_enabled == "Y"),
        ("CloudN / Edge Enabled",      data.cloudn_edge == "Y"),
        ("Multi-Region Transit",       data.multi_region == "Y"),
    ]
    use_cases = [
        ("Unified Cloud NW Fabric",    data.uc_unified_nw == "Y"),
        ("Prevent Lateral Movement",   data.uc_lateral == "Y"),
        ("Zero Trust NW Segmentation", data.uc_zt_seg == "Y"),
        ("Secure 3rd Party Access",    data.uc_3rd_party == "Y"),
        ("End-to-End Encryption",      data.uc_e2e_enc == "Y"),
        ("Block Data Exfiltration",    data.uc_block_exfil == "Y"),
        ("Secure Dev Velocity",        data.uc_dev_velocity == "Y"),
    ]

    LEFT      = 0.529
    TOP       = 1.685
    GAP       = 0.3
    TOTAL_W   = 12.2
    TBL_W     = (TOTAL_W - GAP) / 2   # ~5.95"
    STATUS_W  = 0.65
    NAME_W    = TBL_W - STATUS_W
    HEIGHT    = 5.2

    def _style_cell(cell, text, *, bold=False, size=14, fg=None, bg=None,
                    align=PP_ALIGN.LEFT):
        tf = cell.text_frame
        tf.word_wrap = False
        para = tf.paragraphs[0]
        para.alignment = align
        # Remove existing runs and set fresh text
        for r in para.runs:
            r._r.getparent().remove(r._r)
        run = para.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        if fg:
            run.font.color.rgb = fg
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg

    def _add_panel(items, title, left_in):
        n_rows = 1 + len(items)
        gf  = slide.shapes.add_table(
            n_rows, 2,
            Inches(left_in), Inches(TOP),
            Inches(TBL_W), Inches(HEIGHT),
        )
        tbl = gf.table
        tbl.columns[0].width = Inches(STATUS_W)
        tbl.columns[1].width = Inches(NAME_W)

        # Header: merge both cells, navy background
        hdr = tbl.cell(0, 0)
        hdr.merge(tbl.cell(0, 1))
        _style_cell(hdr, title, bold=True, size=16,
                    fg=_WHITE, bg=_NAVY, align=PP_ALIGN.CENTER)

        # Data rows
        for i, (name, active) in enumerate(items):
            bg  = _ROW_BG[i % 2]
            sym = "✓" if active else "✗"
            col = _GREEN if active else _RED
            _style_cell(tbl.cell(i + 1, 0), sym, bold=True, size=16,
                        fg=col, bg=bg, align=PP_ALIGN.CENTER)
            _style_cell(tbl.cell(i + 1, 1), name, fg=_NAVY, bg=bg)

    _add_panel(features,  "Features Enabled", LEFT)
    _add_panel(use_cases, "Use Cases",         LEFT + TBL_W + GAP)


def update_slide9_feature_requests(slide, data: CustomerData,
                                   fr_tbl_img: Optional[Path]):
    """
    Slide 10 (index 9) – Feature Requests.
    Populates the counts text placeholder and embeds a table image with
    the full feature request details (Key, Summary, Issue Type, Status).
    """
    # ── Counts text placeholder ──────────────────────────────────────────── #
    content_shape = find_shape(slide, "Text Placeholder 8")
    if content_shape is None:
        for sh in slide.shapes:
            if sh.has_text_frame and ("Feature" in sh.text_frame.text or
                                      "Request" in sh.text_frame.text):
                content_shape = sh
                break

    if content_shape is not None and content_shape.has_text_frame:
        tf = content_shape.text_frame
        tf.word_wrap = True
        for p in tf.paragraphs[1:]:
            p._p.getparent().remove(p._p)
        p0 = tf.paragraphs[0]
        if data.has_feature_requests:
            lines = [
                f"Total: {data.fr_total_count}   "
                f"Completed: {data.fr_completed_count}   "
                f"Outstanding: {data.fr_outstanding_count}",
            ]
        else:
            lines = ["No feature requests on record."]
        if p0.runs:
            p0.runs[0].text = lines[0]
            for r in p0.runs[1:]:
                r.text = ""
        else:
            p0.add_run().text = lines[0]
        for line in lines[1:]:
            np_ = tf.add_paragraph()
            np_.add_run().text = line
    else:
        log.warning("Could not find Feature Requests placeholder on slide 9")

    # ── Feature requests detail table image ──────────────────────────────── #
    if data.has_feature_requests:
        LEFT = 0.529
        TOP  = 2.05
        W    = 12.2
        H    = _h(W, cb.FR_FW, cb.FR_FH)
        add_image_inches(slide, fr_tbl_img, LEFT, TOP, W, H)


def update_slide10_strategic(slide, data: CustomerData):
    """
    Slide 11 (index 10) – Strategic Engagements.

    Template has two side-by-side text boxes (both named 'Text Placeholder 8'):
      Left:  (0.529", 1.659") 5.525"×4.591"  – 'Current Projects (PS)\\nNotes'
      Right: (7.429", 1.659") 5.525"×4.591"  – 'Training & Certification\\nNotes'
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        tf   = shape.text_frame

        if "Current Projects (PS)" in text:
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            for b in _ps_ta_bullets(data):
                np_ = tf.add_paragraph()
                np_.level = 1
                np_.add_run().text = b

        elif "Training" in text and "ertif" in text:
            for p in tf.paragraphs[1:]:
                p._p.getparent().remove(p._p)
            for b in _cert_bullets(data):
                np_ = tf.add_paragraph()
                np_.level = 1
                np_.add_run().text = b


def _ps_ta_bullets(data: CustomerData) -> list:
    bullets = []
    if data.ta_engaged:
        bullets.append("TA engagement in progress.")
    else:
        bullets.append("No active TA engagement.")
    if data.ps_hours > 0:
        bullets.append(f"{data.ps_hours} PS hours tracked.")
    return bullets or ["No PS activity recorded."]


def _cert_bullets(data: CustomerData) -> list:
    if data.active_ace_certs > 0:
        return [f"{data.active_ace_certs} active ACE certification(s)."]
    if data.last_ace_date and data.last_ace_date not in ("nan", "NaT", ""):
        out = ["ACE certification(s) expired."]
        if data.next_ace_expiry and data.next_ace_expiry not in ("nan", "NaT", ""):
            out.append(f"Next expiry: {data.next_ace_expiry}.")
        return out
    return ["No certifications on record."]


def update_slide12_score(slide, data: CustomerData):
    """
    Slide 13 (index 12) – Client Success Score.
    Populate 'Text Placeholder 8' with suggested scores.
    OLE 'Object 1' is NOT touched.
    """
    rationale = {
        "Product Utilization":   _pu_rationale(data),
        "Reliability & Support": _rs_rationale(data),
        "Use Case Adoption":     _uca_rationale(data),
        "Strategic Engagement":  _se_rationale(data),
    }

    ph = find_shape(slide, "Text Placeholder 8")
    if ph is None:
        # Fall back: find any notes-like box below the title area
        ph = _find_notes_shape(slide, keywords=["notes"])
    if ph is None or not ph.has_text_frame:
        log.warning("Could not find score text placeholder on slide 12")
        return

    tf = ph.text_frame
    tf.word_wrap = True

    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    header = "Suggested scores. Reference scoring guide."
    if p0.runs:
        p0.runs[0].text = header
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.add_run().text = header

    def _line(text, level=0):
        np_ = tf.add_paragraph()
        np_.level = level
        np_.add_run().text = text

    for pillar, score in data.scores.items():
        label = data.score_label(score)
        _line(f"{pillar}: Score {score} \u2013 {label}")
        _line(rationale.get(pillar, ""), level=1)


def _pu_rationale(data: CustomerData) -> str:
    if data.unsupported_controllers > 0:
        return f"{data.unsupported_controllers} controller(s) on unsupported versions"
    if data.copilot_deployed > 0 and data.bv_status not in ("Not Initiated", ""):
        return "CoPilot deployed; BV assessment completed"
    return "Controllers supported; CoPilot not yet deployed"


def _rs_rationale(data: CustomerData) -> str:
    if data.p1_tickets > 0:
        return f"{data.p1_tickets} P1 severity ticket(s) raised"
    if data.defect_tickets > 0:
        return f"{data.defect_tickets} defect ticket(s), no P1s"
    if data.total_tickets > 0:
        return f"{data.total_tickets} total tickets; no P1/defects"
    return "No tickets in past 12 months"


def _uca_rationale(data: CustomerData) -> str:
    enabled = []
    if data.uc_unified_nw == "Y":   enabled.append("Unified NW")
    if data.uc_zt_seg == "Y":       enabled.append("ZT Seg")
    if data.uc_lateral == "Y":      enabled.append("Lateral Mvmt")
    if data.uc_block_exfil == "Y":  enabled.append("Data Exfil")
    if data.uc_3rd_party == "Y":    enabled.append("3rd Party Access")
    if data.uc_e2e_enc == "Y":      enabled.append("E2E Enc")
    return ("Active: " + ", ".join(enabled)) if enabled else "No Zero Trust use cases active"


def _se_rationale(data: CustomerData) -> str:
    parts = []
    if data.active_ace_certs > 0:
        parts.append(f"{data.active_ace_certs} ACE cert(s)")
    if data.ps_hours > 0:
        parts.append(f"{data.ps_hours} PS hrs")
    if data.ta_engaged:
        parts.append("TA engaged")
    return ", ".join(parts) if parts else "No strategic activity recorded"


def update_slide14_next_steps(slide, data: CustomerData):
    """Slide 15 (index 14) – Next Steps."""
    if data.copilot_deployed == 0:
        for shape in slide.shapes:
            if (shape.has_text_frame and
                    "collaborative design session" in shape.text_frame.text.lower()):
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = (
                        "Enable CoPilot for full network visibility and control"
                    )
                break


# ============================================================
# Main orchestrator
# ============================================================

def build_presentation(data: CustomerData, template_path, output_path: str):
    """
    Generate the CBR PPTX from a pre-built CustomerData object.

    Args:
        data:          Populated CustomerData instance.
        template_path: Path (str/Path) OR bytes-like object for the template PPTX.
        output_path:   Destination file path for the generated PPTX.
    """
    output_path = Path(output_path)
    date_str = datetime.today().strftime("%B %d, %Y")

    with tempfile.TemporaryDirectory() as tmp_str:
        tmp = Path(tmp_str)
        print("Generating chart images …")

        # ── Slide 3: Product Consumption ──────────────────────────────── #
        pivot_mrr, pivot_usage = data.get_consumption_chart_data()

        # Normalise: CSVs may store this as a ratio (0.875) or already as % (87.5)
        _pct_raw = data.consumption_pct
        _pct_display = round(_pct_raw if _pct_raw > 1 else _pct_raw * 100, 1)
        s3_util  = cb.metric_tile(
            f"{_pct_display}%",
            "Consumption Utilization",
            tmp / "s3_util.png",
        )
        s3_stack = cb.stacked_parameter_chart(pivot_mrr, pivot_usage,
                                               tmp / "s3_stack.png")
        s3_spend = cb.monthly_spend_chart(data.df_monthly_spend,
                                           tmp / "s3_spend.png")

        # ── Slide 4: Operational Health ───────────────────────────────── #
        s4_t1 = cb.metric_tile(str(data.supported_controllers),
                                "Supported Controllers",     tmp / "s4_t1.png")
        s4_t2 = cb.metric_tile(str(data.unsupported_controllers),
                                "Unsupported Controllers",   tmp / "s4_t2.png")
        s4_t3 = cb.metric_tile(str(data.copilot_deployed),
                                "CoPilot Enabled Controllers", tmp / "s4_t3.png")
        s4_tbl = cb.controller_table_image(data,            tmp / "s4_tbl.png")

        # ── Slide 6: Support Requests ─────────────────────────────────── #
        s6_t1 = cb.metric_tile(str(data.total_tickets),
                                "Total Tickets",                       tmp / "s6_t1.png")
        s6_t2 = cb.metric_tile(str(data.p1_tickets),
                                "Aviatrix Defects (P1's)",             tmp / "s6_t2.png")
        s6_t3 = cb.metric_tile(str(data.defect_tickets),
                                "Aviatrix Defects (excl. P1's)",       tmp / "s6_t3.png")
        s6_t4 = cb.metric_tile(str(data.other_tickets),
                                "Others / Non-Aviatrix",               tmp / "s6_t4.png")
        s6_req = cb.ticket_requestor_table(data.df_ticket_requestor,   tmp / "s6_req.png")

        # ── Slide 7: Software Upgrades ────────────────────────────────── #
        s7_t1  = cb.metric_tile(str(data.upgrade_tickets_count),
                                 "Product Upgrade Tickets",  tmp / "s7_t1.png")
        s7_t2  = cb.metric_tile(str(data.gateway_sum),
                                 "Total Gateways",           tmp / "s7_t2.png")
        s7_tbl = cb.release_lifecycle_table(data,            tmp / "s7_tbl.png")

        # Slide 8 (Use Cases) uses native PPTX tables — no image needed

        # ── Slide 9: Feature Requests ─────────────────────────────────── #
        s9_fr = (cb.feature_requests_table(data.df_feature_requests, tmp / "s9_fr.png")
                 if data.has_feature_requests else None)

        # ── Manual slide images ───────────────────────────────────────── #
        manual_full   = cb.manual_slide_placeholder(tmp / "manual_full.png",  short=False)
        manual_banner = cb.manual_slide_placeholder(tmp / "manual_banner.png", short=True)

        # ── Build presentation ────────────────────────────────────────── #
        print("Building presentation …")
        # template_path may be a path string/Path, bytes, or BytesIO object
        import io as _io
        if isinstance(template_path, (bytes, bytearray)):
            prs = Presentation(_io.BytesIO(template_path))
        elif isinstance(template_path, _io.IOBase):
            prs = Presentation(template_path)
        else:
            prs = Presentation(str(template_path))

        n_slides = len(prs.slides)
        if n_slides < 16:
            raise ValueError(
                f"Template has only {n_slides} slide(s); expected at least 16. "
                "Ensure you are using the correct CBR Template.pptx."
            )
        slides = prs.slides

        update_slide0_title(slides[0], data, date_str)
        update_slide1_agenda(slides[1], data)

        _mark_manual(slides[2], data, manual_full)              # Customer Priorities

        update_slide3_consumption(slides[3], data,
                                   s3_util, s3_stack, s3_spend)
        update_slide4_operational(slides[4], data,
                                   s4_t1, s4_t2, s4_t3, s4_tbl)
        update_slide5_business_value(slides[5], data)
        update_slide6_reliability(slides[6], data,
                                   s6_t1, s6_t2, s6_t3, s6_t4, s6_req)
        update_slide7_software_upgrades(slides[7], data,
                                         s7_t1, s7_t2, s7_tbl)
        update_slide8_usecases(slides[8], data)                   # Use Cases
        update_slide9_feature_requests(slides[9], data, s9_fr)    # Feature Requests
        update_slide10_strategic(slides[10], data)

        _mark_manual(slides[11], data, manual_banner, short=True)  # ZTMM (Picture 10 preserved)

        update_slide12_score(slides[12], data)

        _mark_manual(slides[13], data, manual_full)              # Current Initiatives
        update_slide14_next_steps(slides[14], data)
        _mark_manual(slides[15], data, manual_full)              # Action Plan

        # ── Delete rubric slides (highest index first) ────────────────── #
        indices_to_delete = sorted(
            [i - 1 for i in SLIDES_TO_REMOVE_1IDX], reverse=True
        )
        print(f"Removing {len(indices_to_delete)} rubric slides …")
        for idx in indices_to_delete:
            if idx < len(prs.slides):
                delete_slide(prs, idx)

        prs.save(str(output_path))
        print(f"\n✓  Saved: {output_path}  ({len(prs.slides)} slides)")


# ============================================================
# CLI wrapper (kept for local use)
# ============================================================

def generate_cbr(data_dir: str, template_path: str, output_path: str):
    """CLI-friendly wrapper: loads CustomerData then calls build_presentation."""
    data_dir      = Path(data_dir)
    template_path = Path(template_path)

    if not data_dir.exists():
        raise FileNotFoundError(f"Data directory not found: {data_dir}")
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    print(f"Loading customer data from: {data_dir}")
    data = CustomerData(data_dir)
    print(f"  Customer : {data.customer_name}")
    print(f"  ARR      : ${data.arr:,.0f}")
    print(f"  Scores   : {data.scores}")

    build_presentation(data, template_path, output_path)


def main():
    parser = argparse.ArgumentParser(
        description="Generate an Aviatrix Customer Business Review deck."
    )
    parser.add_argument(
        "--data-dir", default=DEFAULT_DATA_DIR,
        help=f"CSV data folder  (default: {DEFAULT_DATA_DIR})",
    )
    parser.add_argument(
        "--template", default=DEFAULT_TEMPLATE,
        help=f"PPTX template    (default: {DEFAULT_TEMPLATE})",
    )
    parser.add_argument(
        "--output", default=None,
        help="Output filename   (auto-generated from customer name if omitted)",
    )
    args = parser.parse_args()

    try:
        data_dir      = Path(args.data_dir)
        template_path = Path(args.template)
        if not data_dir.exists():
            raise FileNotFoundError(f"Data directory not found: {data_dir}")
        if not template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        data = CustomerData(data_dir)
        log.info("Customer: %s  ARR: $%,.0f  Scores: %s",
                 data.customer_name, data.arr, data.scores)

        if args.output:
            output = args.output
        else:
            safe   = data.customer_name.replace(" ", "_").replace("/", "-")
            slug   = datetime.today().strftime("%Y-%m-%d")
            output = f"CBR - {safe} - {slug}.pptx"

        build_presentation(data, template_path, output)
    except FileNotFoundError as e:
        print(f"ERROR: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
